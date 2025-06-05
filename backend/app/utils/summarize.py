import fitz  # PyMuPDF for PDF
import pptx  # python-pptx for PowerPoint
import docx  # python-docx for Word documents
import google.generativeai as genai
import os
from typing import List, Dict, Tuple, Optional
from dataclasses import dataclass
import re
import time
import traceback

@dataclass
class DocumentSection:
    title: str
    content: str
    importance: float = 1.0
    depth: int = 1  # Hierarchy depth (1 = main section, 2 = subsection, etc.)

class QuotaFriendlyAnalyzer:
    def __init__(self, api_key: str):
        self.api_key = api_key
        genai.configure(api_key=api_key)
        # Use the latest Gemini model
        self.model = genai.GenerativeModel('gemini-2.0-flash')
        self.api_calls = 0  # Track number of API calls
        
    def extract_pdf_text(self, pdf_path: str) -> List[DocumentSection]:
        """Extract text from PDF with improved section detection"""
        try:
            pdf_document = fitz.open(pdf_path)
            sections = []
            current_text = []
            current_title = "Introduction"
            current_depth = 1
            
            # First pass: analyze document structure to identify fonts and styles
            font_sizes = []
            for page_num in range(min(5, pdf_document.page_count)):  # Sample first 5 pages
                page = pdf_document.load_page(page_num)
                blocks = page.get_text("dict")["blocks"]
                
                for block in blocks:
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line["spans"]:
                                if span["size"] > 0:  # Valid font size
                                    font_sizes.append(span["size"])
            
            # Calculate font size thresholds for headers
            font_sizes.sort()
            if len(font_sizes) > 10:  # Enough samples
                # Use percentiles to determine heading levels
                heading1_threshold = font_sizes[int(len(font_sizes) * 0.9)]
                heading2_threshold = font_sizes[int(len(font_sizes) * 0.8)]
                heading3_threshold = font_sizes[int(len(font_sizes) * 0.7)]
            else:
                # Default thresholds if not enough samples
                heading1_threshold = 16
                heading2_threshold = 14
                heading3_threshold = 12
            
            # Second pass: extract text with section detection
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                blocks = page.get_text("dict")["blocks"]
                
                for block in blocks:
                    if "lines" in block:
                        formatted_text = []
                        for line in block["lines"]:
                            for span in line["spans"]:
                                text = span["text"].strip()
                                size = span["size"]
                                flags = span["flags"]
                                
                                # Detect section headers with more granularity
                                is_header = False
                                new_depth = 0
                                
                                if size >= heading1_threshold or (flags & 16 and len(text) < 100 and text.strip()):
                                    is_header = True
                                    new_depth = 1  # Main heading
                                elif size >= heading2_threshold or (flags & 4 and len(text) < 100 and text.strip()):
                                    is_header = True
                                    new_depth = 2  # Subheading
                                elif size >= heading3_threshold and len(text) < 100 and text.strip():
                                    is_header = True
                                    new_depth = 3  # Sub-subheading
                                
                                if is_header and text and not text.isdigit() and len(text) > 1:
                                    # Save current section before starting new one
                                    if current_text:
                                        sections.append(DocumentSection(
                                            title=current_title,
                                            content="\n".join(current_text),
                                            depth=current_depth
                                        ))
                                    current_title = text
                                    current_text = []
                                    current_depth = new_depth
                                elif text:
                                    formatted_text.append(text)
                        
                        if formatted_text:
                            current_text.append(" ".join(formatted_text))
            
            # Add final section
            if current_text:
                sections.append(DocumentSection(
                    title=current_title,
                    content="\n".join(current_text),
                    depth=current_depth
                ))
            
            # Post-process to combine very small sections with next section
            i = 0
            while i < len(sections) - 1:
                if len(sections[i].content.split()) < 30:  # Less than 30 words
                    # Combine with next section
                    sections[i+1].content = f"{sections[i].title}:\n{sections[i].content}\n\n{sections[i+1].content}"
                    sections.pop(i)
                else:
                    i += 1
            
            pdf_document.close()
            return sections
        
        except Exception as e:
            print(f"PDF extraction error: {e}")
            return []

    def extract_pptx_text(self, pptx_path: str) -> List[DocumentSection]:
        """Extract text from PowerPoint with improved structure detection"""
        try:
            presentation = pptx.Presentation(pptx_path)
            sections = []
            
            # First determine possible slide types and structure
            has_section_slides = False
            
            for slide in presentation.slides:
                # Check if this looks like a section/divider slide
                if slide.shapes.title and len(slide.shapes) < 3:
                    title = slide.shapes.title.text
                    if len(title) < 50 and any(keyword in title.lower() for keyword in ['section', 'part', 'chapter', 'agenda']):
                        has_section_slides = True
            
            # Now extract content with awareness of structure
            current_section = "Introduction"
            current_slides = []
            
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                slide_title = f"Slide {i+1}"
                
                # Extract title if exists
                if slide.shapes.title:
                    slide_title = slide.shapes.title.text
                
                # Check if this is a section slide
                is_section_slide = False
                if has_section_slides and slide.shapes.title and len(slide.shapes) < 3:
                    is_section_slide = True
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Skip title text as we already captured it
                        if shape != slide.shapes.title:
                            slide_text.append(shape.text)
                
                # Process based on slide type
                if is_section_slide:
                    # Save previous section
                    if current_slides:
                        combined_content = "\n\n".join(current_slides)
                        sections.append(DocumentSection(
                            title=current_section,
                            content=combined_content,
                            depth=1
                        ))
                    
                    # Start new section
                    current_section = slide_title
                    current_slides = []
                else:
                    # Format slide content
                    formatted_content = f"### {slide_title}\n" + "\n".join(slide_text)
                    current_slides.append(formatted_content)
            
            # Add final section
            if current_slides:
                combined_content = "\n\n".join(current_slides)
                sections.append(DocumentSection(
                    title=current_section,
                    content=combined_content,
                    depth=1
                ))
            
            # Handle case where no sections were detected
            if not sections and has_section_slides == False:
                all_slides = []
                for i, slide in enumerate(presentation.slides):
                    slide_text = []
                    slide_title = f"Slide {i+1}"
                    
                    if slide.shapes.title:
                        slide_title = slide.shapes.title.text
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                            slide_text.append(shape.text)
                    
                    formatted_content = f"### {slide_title}\n" + "\n".join(slide_text)
                    all_slides.append(formatted_content)
                
                sections.append(DocumentSection(
                    title="Presentation Content",
                    content="\n\n".join(all_slides),
                    depth=1
                ))
            
            return sections
        
        except Exception as e:
            print(f"PowerPoint extraction error: {e}")
            return []

    def extract_docx_text(self, docx_path: str) -> List[DocumentSection]:
        """Extract text from Word document with improved structure recognition"""
        try:
            doc = docx.Document(docx_path)
            sections = []
            current_section_title = "Introduction"
            current_section_text = []
            current_depth = 1
            
            # Map Word heading styles to depth levels
            heading_styles = {
                'Title': 0,
                'Heading 1': 1, 
                'Heading 2': 2, 
                'Heading 3': 3,
                'Heading 4': 4
            }
            
            # Process paragraphs
            for paragraph in doc.paragraphs:
                # Skip empty paragraphs
                if not paragraph.text.strip():
                    continue
                
                # Detect headers based on style
                is_header = False
                new_depth = 0
                
                if paragraph.style.name in heading_styles:
                    is_header = True
                    new_depth = heading_styles[paragraph.style.name]
                # Fallback detection based on formatting
                elif len(paragraph.text) < 100:  # Potential header by length
                    if paragraph.style.font.bold or paragraph.style.font.size and paragraph.style.font.size >= 14:
                        is_header = True
                        new_depth = 1 if paragraph.style.font.size and paragraph.style.font.size >= 16 else 2
                
                if is_header:
                    # Save previous section if not empty
                    if current_section_text:
                        sections.append(DocumentSection(
                            title=current_section_title,
                            content="\n".join(current_section_text),
                            depth=current_depth
                        ))
                    
                    # Start new section
                    current_section_title = paragraph.text
                    current_section_text = []
                    current_depth = new_depth
                else:
                    # Add paragraph text to current section
                    current_section_text.append(paragraph.text)
            
            # Add final section
            if current_section_text:
                sections.append(DocumentSection(
                    title=current_section_title,
                    content="\n".join(current_section_text),
                    depth=current_depth
                ))
            
            # Post-process: merge very small sections
            i = 0
            while i < len(sections) - 1:
                if len(sections[i].content.split()) < 30:  # Less than 30 words
                    sections[i+1].content = f"{sections[i].title}:\n{sections[i].content}\n\n{sections[i+1].content}"
                    sections.pop(i)
                else:
                    i += 1
            
            return sections
        
        except Exception as e:
            print(f"Word document extraction error: {e}")
            return []

    def calculate_section_importance(self, sections: List[DocumentSection]) -> List[DocumentSection]:
        """Calculate importance score for each section based on content and position"""
        for i, section in enumerate(sections):
            # Base importance on position (higher for earlier sections)
            position_weight = max(1.0, 1.5 - (i / len(sections)))
            
            # Content length factor (longer sections might be more important)
            length_weight = min(1.5, max(0.5, len(section.content.split()) / 500))
            
            # Title importance (introductions, conclusions, etc.)
            title_weight = 1.0
            important_keywords = ['introduction', 'conclusion', 'summary', 'result', 'finding', 'discussion', 'analysis']
            if any(keyword in section.title.lower() for keyword in important_keywords):
                title_weight = 1.5
            
            # Depth weight (main sections are more important than subsections)
            depth_weight = 1.5 / max(1, section.depth)
            
            # Calculate final importance score
            section.importance = position_weight * length_weight * title_weight * depth_weight
        
        return sections

    def prepare_batch_content(self, sections: List[DocumentSection], max_batch_size: int = 10000) -> List[str]:
        """Prepare document content in batches for efficient API usage"""
        batches = []
        current_batch = []
        current_size = 0
        
        for section in sections:
            section_text = f"## {section.title}\n{section.content}"
            section_size = len(section_text)
            
            # If this section alone exceeds max size, split it
            if section_size > max_batch_size:
                # Add current batch if not empty
                if current_batch:
                    batches.append("\n\n".join(current_batch))
                    current_batch = []
                    current_size = 0
                
                # Split large section
                words = section_text.split()
                chunk_size = max_batch_size // 2  # characters per chunk (approximate)
                current_chunk = []
                current_chunk_size = 0
                
                for word in words:
                    word_size = len(word) + 1  # +1 for space
                    if current_chunk_size + word_size > chunk_size and current_chunk:
                        # Complete this chunk
                        chunk_text = f"## {section.title} (Part {len(batches) + 1})\n" + " ".join(current_chunk)
                        batches.append(chunk_text)
                        current_chunk = [word]
                        current_chunk_size = word_size
                    else:
                        current_chunk.append(word)
                        current_chunk_size += word_size
                
                # Add final chunk
                if current_chunk:
                    chunk_text = f"## {section.title} (Part {len(batches) + 1})\n" + " ".join(current_chunk)
                    batches.append(chunk_text)
            
            # Normal case: try to add to current batch
            elif current_size + section_size <= max_batch_size:
                current_batch.append(section_text)
                current_size += section_size
            else:
                # Finish current batch and start new one
                batches.append("\n\n".join(current_batch))
                current_batch = [section_text]
                current_size = section_size
        
        # Add final batch if not empty
        if current_batch:
            batches.append("\n\n".join(current_batch))
        
        return batches

    def summarize_batch(self, batch_content: str, file_name: str, batch_index: int, total_batches: int) -> str:
        """Generate a comprehensive summary for a batch of sections with rate limit awareness"""
        try:
            # Create an educational and informative prompt
            batch_description = "full document" if total_batches == 1 else f"batch {batch_index+1} of {total_batches}"
            
            prompt = f"""I need a highly educational and informative summary of the following document content from {file_name} ({batch_description}).

            YOUR TASK:
            Create an advanced yet accessible summary that would help someone with no prior knowledge understand this topic deeply. Your summary should:

            1. Begin with a concise overview of what this content covers
            2. Explain complex concepts in clear, educational language
            3. Identify and elaborate on key points, insights, or arguments
            4. Define any specialized terminology or jargon
            5. Explain real-world implications or applications when relevant
            6. Use markdown formatting for readability (headings, bullet points, etc.)

            Make your language clear but sophisticated, as if writing for an educated reader trying to master a new subject.

            DOCUMENT CONTENT:
            {batch_content}
            """
            
            # Track API calls and add simple rate limiting
            self.api_calls += 1
            if self.api_calls > 1:
                # Add delay between calls to avoid rate limiting
                time.sleep(2)
            
            # Generate summary with simple retry mechanism
            max_retries = 2
            for attempt in range(max_retries):
                try:
                    response = self.model.generate_content(
                        prompt,
                        generation_config={
                            'temperature': 0.2,
                            'top_p': 0.95,
                            'max_output_tokens': 4000
                        }
                    )
                    
                    # Check if we got a valid response
                    if response.text and len(response.text) > 100:
                        return response.text
                    
                    # If response is too short, retry with higher temperature
                    time.sleep(3)  # Longer pause before retry
                    
                except Exception as inner_e:
                    print(f"Generation error (attempt {attempt+1}): {inner_e}")
                    # Only sleep if we're going to retry
                    if attempt < max_retries - 1:
                        time.sleep(5)  # Much longer pause after error
            
            # Fallback if all retries failed
            return f"""# Summary of {file_name} ({batch_description})

            Unfortunately, I wasn't able to generate a detailed summary for this content due to API limitations. Here's a basic overview:

            This section appears to cover various concepts related to the document's subject matter. The content includes several key sections and technical information that would be valuable for understanding the topic.

            **Note:** A more detailed summary could not be generated at this time. Please try again later or contact support if this issue persists.
            """
            
        except Exception as e:
            print(f"Batch summarization error: {e}")
            return f"Error summarizing batch: {str(e)}"

    def create_comprehensive_summary(self, file_path: str) -> str:
        """Create an advanced, educational summary of the document with quota awareness"""
        try:
            # Reset API call counter
            self.api_calls = 0
            
            # Get file name and extension
            file_name = os.path.basename(file_path)
            file_extension = os.path.splitext(file_path)[1].lower()
            
            # Extract sections based on file type
            if file_extension == '.pdf':
                sections = self.extract_pdf_text(file_path)
            elif file_extension == '.pptx':
                sections = self.extract_pptx_text(file_path)
            elif file_extension == '.docx':
                sections = self.extract_docx_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")
            
            if not sections:
                return "No content could be extracted from the document."
            
            # Calculate section importance
            sections = self.calculate_section_importance(sections)
            
            # Sort sections by importance for better analysis
            sections.sort(key=lambda x: x.importance, reverse=True)
            
            # Take top 80% of sections by importance to focus on key content
            cutoff_index = int(len(sections) * 0.8)
            selected_sections = sections[:max(3, cutoff_index)]  # At least 3 sections
            
            # Sort selected sections back to original order
            selected_sections.sort(key=lambda x: sections.index(x))
            
            # Prepare content in batches to avoid token limits
            content_batches = self.prepare_batch_content(selected_sections, max_batch_size=12000)
            
            # Generate summary for each batch (this limits API calls)
            batch_summaries = []
            for i, batch in enumerate(content_batches):
                # Only process up to 3 batches maximum to avoid excessive API calls
                if i >= 3:
                    batch_summaries.append(f"# Additional Content\n\nThis document contains more content that wasn't summarized due to length constraints.")
                    break
                    
                summary = self.summarize_batch(batch, file_name, i, len(content_batches))
                batch_summaries.append(summary)
            
            # Combine all summaries
            if len(batch_summaries) == 1:
                final_output = f"""# Advanced Educational Summary: {file_name}

{batch_summaries[0]}

---

*This summary was generated automatically and presents key concepts in an educational format.*
"""
            else:
                # For multiple batches, add a table of contents
                toc = "## Table of Contents\n\n"
                content = ""
                
                for i, summary in enumerate(batch_summaries):
                    section_title = f"Part {i+1}: {file_name}"
                    toc += f"{i+1}. [{section_title}](#{section_title.lower().replace(' ', '-')})\n"
                    content += f"# {section_title}\n\n{summary}\n\n---\n\n"
                
                final_output = f"""# Advanced Educational Summary: {file_name}

{toc}

{content}

*This summary was generated automatically and presents key concepts in an educational format.*
"""
            
            return final_output
            
        except Exception as e:
            print(f"An error occurred during analysis: {str(e)}")
            print(traceback.format_exc())
            return f"An error occurred during analysis: {str(e)}"