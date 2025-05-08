import fitz  # PyMuPDF
import google.generativeai as genai
import os
import json
import re
from pptx import Presentation
import pdfplumber
from typing import List, Dict

class FlashcardGenerator:
    def __init__(self, api_key: str, model: str = "gemini-1.5-pro"):
        """
        Initialize the FlashcardGenerator with Gemini API.
        :param api_key: Google Gemini API key
        :param model: Gemini model to use (default: gemini-1.5-pro)
        """
        genai.configure(api_key=api_key)
        self.model = genai.GenerativeModel(model)
        self.generation_config = {
            "temperature": 0.2,
            "top_p": 0.95,
            "top_k": 40,
            "max_output_tokens": 8192,
        }
        
    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF using multiple methods for better reliability"""
        text = ""
        # First try with PyMuPDF
        try:
            pdf_document = fitz.open(file_path)
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                text += page.get_text("text") + "\n\n"
            pdf_document.close()
        except Exception as e:
            print(f"PyMuPDF extraction failed: {e}")
            
        # If PyMuPDF didn't get much text, try pdfplumber
        if len(text.strip()) < 100:
            try:
                text = ""
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        extracted = page.extract_text()
                        if extracted:
                            text += extracted + "\n\n"
            except Exception as e:
                print(f"pdfplumber extraction failed: {e}")
                
        return text.strip()

    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations"""
        try:
            presentation = Presentation(file_path)
            text = []
            
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                
                # Extract title if present
                if slide.shapes.title and slide.shapes.title.text:
                    slide_text.append(f"Slide {i+1} Title: {slide.shapes.title.text}")
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        # Avoid duplicating the title
                        if not (slide.shapes.title and shape.text == slide.shapes.title.text):
                            slide_text.append(shape.text)
                
                # Add slide content to overall text
                if slide_text:
                    text.append("\n".join(slide_text))
            
            return "\n\n".join(text)
        except Exception as e:
            print(f"Error extracting text from PPTX: {e}")
            return ""

    def extract_text_from_file(self, file_path: str) -> str:
        """Extract text from supported file types"""
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        if ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self.extract_text_from_pptx(file_path)
        elif ext in ['.txt', '.md', '.csv']:
            # Simple text files
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except UnicodeDecodeError:
                # Try another encoding if utf-8 fails
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
        else:
            print(f"Unsupported file type: {ext}")
            return ""

    def estimate_tokens(self, text: str) -> int:
        """Estimate the number of tokens in text (rough approximation)"""
        # A very rough approximation - assumes 4 chars per token on average
        return len(text) // 4

    def split_text_into_chunks(self, text: str, max_tokens: int = 6000) -> List[str]:
        """
        Split text into chunks that won't exceed Gemini token limits.
        """
        # Make a rough estimate of max characters
        max_chars = max_tokens * 4
        
        # Split text into paragraphs
        paragraphs = text.split('\n\n')
        chunks = []
        current_chunk = []
        current_chars = 0
        
        for paragraph in paragraphs:
            paragraph_chars = len(paragraph)
            
            # If a single paragraph is too large, split it further
            if paragraph_chars > max_chars:
                # If we have anything in the current chunk, add it
                if current_chunk:
                    chunks.append('\n\n'.join(current_chunk))
                    current_chunk = []
                    current_chars = 0
                
                # Split paragraph into sentences
                sentences = re.split(r'(?<=[.!?])\s+', paragraph)
                sentence_chunk = []
                sentence_chars = 0
                
                for sentence in sentences:
                    sentence_length = len(sentence)
                    
                    # If this sentence would make the chunk too large
                    if sentence_chars + sentence_length > max_chars:
                        if sentence_chunk:
                            chunks.append(' '.join(sentence_chunk))
                            sentence_chunk = []
                            sentence_chars = 0
                    
                    # Add the sentence to the current sentence chunk
                    sentence_chunk.append(sentence)
                    sentence_chars += sentence_length
                
                # Add any remaining sentences
                if sentence_chunk:
                    chunks.append(' '.join(sentence_chunk))
                
            # If adding this paragraph would exceed the limit
            elif current_chars + paragraph_chars > max_chars:
                chunks.append('\n\n'.join(current_chunk))
                current_chunk = [paragraph]
                current_chars = paragraph_chars
            # Otherwise, add to current chunk
            else:
                current_chunk.append(paragraph)
                current_chars += paragraph_chars
        
        # Add any remaining content
        if current_chunk:
            chunks.append('\n\n'.join(current_chunk))
            
        return chunks

    def generate_flashcards_for_chunk(self, text: str, num_cards: int = 10) -> List[Dict[str, str]]:
        """Generate flashcards for a text chunk using Gemini API"""
        if not text.strip():
            return []
            
        prompt = f"""Create exactly {num_cards} high-quality flashcards from this text. Focus on important concepts, definitions, and key relationships.
        
        For each flashcard:
        1. Create a specific question that tests understanding
        2. Provide a comprehensive but concise answer
        3. Include numerical data and specific details when relevant
        4. Create questions that promote critical thinking
        5. Ensure questions are diverse (conceptual, factual, and analytical)

        Return ONLY a JSON array with this exact format:
        [
            {{"question": "Question text here?", "answer": "Answer text here."}},
            {{"question": "Another question?", "answer": "Another answer."}}
        ]

        Text to analyze:
        {text}
        """
        
        try:
            response = self.model.generate_content(
                prompt,
                generation_config=self.generation_config
            )
            
            # Extract JSON from response
            response_text = response.text
            
            # Look for JSON array in the response
            json_match = re.search(r'\[\s*{.*}\s*\]', response_text, re.DOTALL)
            if json_match:
                json_str = json_match.group(0)
                try:
                    return json.loads(json_str)
                except json.JSONDecodeError:
                    # Clean up common JSON issues
                    json_str = re.sub(r',\s*}', '}', json_str)  # Remove trailing commas
                    json_str = re.sub(r',\s*]', ']', json_str)  # Remove trailing commas
                    return json.loads(json_str)
            
            # If we didn't find a JSON array pattern, try loading the whole response
            try:
                return json.loads(response_text)
            except json.JSONDecodeError:
                print(f"Failed to parse JSON from response: {response_text[:200]}...")
                return []
                
        except Exception as e:
            print(f"Error generating flashcards: {str(e)}")
            return []

    def generate_flashcards(self, text: str, num_flashcards: int = 10) -> List[Dict[str, str]]:
        """Generate flashcards from a document's text content"""
        if not text or len(text.strip()) < 50:
            print("Text is too short to generate flashcards")
            return []

        # Split text into manageable chunks
        chunks = self.split_text_into_chunks(text)
        
        if not chunks:
            return []
            
        all_flashcards = []
        
        # Calculate how many cards to generate per chunk
        cards_per_chunk = max(2, num_flashcards // len(chunks))
        
        # Generate for each chunk
        for i, chunk in enumerate(chunks):
            # For the last chunk, request any remaining cards
            if i == len(chunks) - 1:
                remaining = num_flashcards - len(all_flashcards)
                if remaining > 0:
                    cards_per_chunk = remaining
            
            chunk_cards = self.generate_flashcards_for_chunk(chunk, cards_per_chunk)
            all_flashcards.extend(chunk_cards)
            
            # If we already have enough, break early
            if len(all_flashcards) >= num_flashcards:
                break
                
        # Remove any duplicate questions
        seen_questions = set()
        unique_flashcards = []
        
        for card in all_flashcards:
            if not isinstance(card, dict) or 'question' not in card or 'answer' not in card:
                continue
                
            # Normalize the question to detect duplicates
            norm_question = card['question'].lower().strip('?!. ')
            if norm_question not in seen_questions:
                seen_questions.add(norm_question)
                unique_flashcards.append(card)
                
        # Trim to requested number
        return unique_flashcards[:num_flashcards]